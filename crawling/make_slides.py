# -*- coding: utf-8 -*-
"""5분 발표용 슬라이드 생성 (python-pptx)"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ── 색상 (Coral Energy 팔레트: 화장품=코랄/골드, AI/기술=네이비) ──
NAVY   = RGBColor(0x2F, 0x3C, 0x7E)
CORAL  = RGBColor(0xF9, 0x61, 0x67)
GOLD   = RGBColor(0xE0, 0xA9, 0x2A)
GREEN  = RGBColor(0x12, 0xA4, 0x7E)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
TEXT   = RGBColor(0x21, 0x29, 0x3D)
GRAY   = RGBColor(0x6B, 0x74, 0x8B)
LIGHT  = RGBColor(0xF4, 0xF5, 0xF9)

FONT = "맑은 고딕"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
W, H = 13.333, 7.5


def add_slide(bg=WHITE):
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = bg
    return s


def textbox(slide, x, y, w, h, text, size=18, color=TEXT, bold=False,
            align=PP_ALIGN.LEFT, italic=False, font=FONT, line_spacing=None,
            anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    lines = text.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if line_spacing:
            p.line_spacing = line_spacing
        r = p.add_run()
        r.text = line
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.italic = italic
        r.font.name = font
        r.font.color.rgb = color
    return tb


def bullets(slide, x, y, w, h, items, size=15, color=TEXT, bold_first=False,
            font=FONT, space_after=8, bullet_color=None):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(space_after)
        r = p.add_run()
        r.text = "•  " + item
        r.font.size = Pt(size)
        r.font.name = font
        r.font.color.rgb = color
    return tb


def rect(slide, x, y, w, h, color, line=False, line_color=None, line_w=1.0, shadow=False):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    if line:
        shp.line.color.rgb = line_color or color
        shp.line.width = Pt(line_w)
    else:
        shp.line.fill.background()
    shp.shadow.inherit = False
    return shp


def rounded_rect(slide, x, y, w, h, color, radius=0.08):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    shp.line.fill.background()
    shp.shadow.inherit = False
    try:
        shp.adjustments[0] = radius
    except Exception:
        pass
    return shp


def oval(slide, x, y, d, color):
    shp = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(d), Inches(d))
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    shp.line.fill.background()
    shp.shadow.inherit = False
    return shp


# ════════════════════════════════════════════════════════════════
# 슬라이드 1 — 표지
# ════════════════════════════════════════════════════════════════
s = add_slide(NAVY)
oval(s, 5.9, 1.0, 1.5, CORAL)
textbox(s, 5.9, 1.0, 1.5, 1.5, "AI", size=32, color=WHITE, bold=True, align=PP_ALIGN.CENTER,
        anchor=MSO_ANCHOR.MIDDLE)

textbox(s, 1.0, 3.0, 11.33, 1.2, "화장품 리뷰 감성분석 서비스",
        size=40, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
textbox(s, 1.0, 4.05, 11.33, 0.6,
        "무신사 · 올리브영 · 쿠팡 리뷰로 학습한 긍정/부정 분류 모델",
        size=18, color=GOLD, align=PP_ALIGN.CENTER, italic=True)
textbox(s, 1.0, 6.6, 11.33, 0.5, "AI 서비스 개발 프로젝트 발표",
        size=13, color=RGBColor(0xC8, 0xCF, 0xE8), align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════════
# 슬라이드 2 — 데이터 수집 과정
# ════════════════════════════════════════════════════════════════
s = add_slide(WHITE)
textbox(s, 0.7, 0.5, 11.9, 0.7, "01. 데이터 수집 과정", size=30, color=NAVY, bold=True)

# 좌측: 수집 현황 표
table_x, table_y, table_w, table_h = 0.7, 1.6, 5.6, 2.6
rows = [
    ("플랫폼", "건수", "방식"),
    ("무신사", "231,029건", "직접 크롤링"),
    ("올리브영", "222,451건", "임현우 크롤링"),
    ("쿠팡", "78,484건", "김동호 크롤링"),
    ("합계", "531,964건", "3명 데이터 통합"),
]
row_h = table_h / len(rows)
col_w = [2.0, 1.6, 2.0]
for ri, row in enumerate(rows):
    is_header = ri == 0
    is_total = ri == len(rows) - 1
    bg = NAVY if is_header else (LIGHT if is_total else WHITE)
    rect(s, table_x, table_y + ri * row_h, table_w, row_h, bg,
         line=True, line_color=RGBColor(0xE2, 0xE5, 0xEE), line_w=0.75)
    cx = table_x
    for ci, cell in enumerate(row):
        color = WHITE if is_header else (NAVY if is_total else TEXT)
        textbox(s, cx + 0.15, table_y + ri * row_h, col_w[ci] - 0.2, row_h, cell,
                size=14 if not is_header else 14, color=color,
                bold=is_header or is_total, anchor=MSO_ANCHOR.MIDDLE)
        cx += col_w[ci]

textbox(s, 0.7, 4.5, 5.6, 0.4, "총 50만 건이 넘는 리뷰를 3개 플랫폼에서 통합 수집",
        size=13, color=GRAY, italic=True)

# 우측: 어려움 카드 2개
card_x, card_w = 6.7, 5.9
rounded_rect(s, card_x, 1.6, card_w, 1.55, LIGHT, radius=0.06)
textbox(s, card_x + 0.3, 1.78, card_w - 0.6, 0.4, "어려움 ① 선택자 오류", size=15, color=CORAL, bold=True)
textbox(s, card_x + 0.3, 2.18, card_w - 0.6, 0.85,
        "React 기반 사이트라 처음 잡은 HTML 선택자가 모두 틀려 상품이 0개 검색됨\n"
        "→ 해결: 개발자도구로 직접 구조 분석, 실제 속성(data-item-id) 적용",
        size=12.5, color=TEXT, line_spacing=1.15)

rounded_rect(s, card_x, 3.35, card_w, 1.55, LIGHT, radius=0.06)
textbox(s, card_x + 0.3, 3.53, card_w - 0.6, 0.4, "어려움 ② 무한 페이지네이션", size=15, color=CORAL, bold=True)
textbox(s, card_x + 0.3, 3.93, card_w - 0.6, 0.85,
        "“이전 후기 보기” 클릭 시 같은 페이지를 무한 반복 크롤링\n"
        "→ 해결: 실제로는 무한스크롤 방식 → 스크롤 시뮬레이션으로 전환",
        size=12.5, color=TEXT, line_spacing=1.15)

# ════════════════════════════════════════════════════════════════
# 슬라이드 3 — 라벨링 문제와 해결 (가장 중요)
# ════════════════════════════════════════════════════════════════
s = add_slide(WHITE)
textbox(s, 0.7, 0.5, 11.9, 0.7, "02. 라벨링 문제와 해결", size=30, color=NAVY, bold=True)
textbox(s, 0.7, 1.15, 11.9, 0.4,
        "긍정/부정 정답이 있어야 학습이 가능한데, 평점만 쓰면 부정 데이터가 거의 없었다",
        size=14, color=GRAY, italic=True)

stat_y, stat_h, stat_w, gap = 1.9, 2.5, 3.6, 0.35
stats = [
    ("0.4%", "무신사 평점 기준", "부정 리뷰가 거의 없음\n(만족도 높은 구매자만 평점)", CORAL),
    ("21.9%", "키워드 추측 (1차 시도)", "부정 키워드로 직접 분류\n→ 오탐 다수 발생", GOLD),
    ("25.0%", "3사 통합 + 실제 평점 (최종)", "올리브영·쿠팡은 부정 비율\n8~17%로 충분 → 통합 후 보정", GREEN),
]
total_w = stat_w * 3 + gap * 2
start_x = (W - total_w) / 2
for i, (num, label, desc, color) in enumerate(stats):
    x = start_x + i * (stat_w + gap)
    rounded_rect(s, x, stat_y, stat_w, stat_h, LIGHT, radius=0.05)
    rect(s, x, stat_y, 0.12, stat_h, color)
    textbox(s, x + 0.35, stat_y + 0.25, stat_w - 0.6, 0.9, num, size=44, color=color, bold=True)
    textbox(s, x + 0.35, stat_y + 1.1, stat_w - 0.6, 0.4, label, size=13.5, color=NAVY, bold=True)
    textbox(s, x + 0.35, stat_y + 1.55, stat_w - 0.6, 0.85, desc, size=11.5, color=TEXT, line_spacing=1.2)

textbox(s, 0.7, 4.75, 11.9, 0.4, "최종 해결 방법", size=16, color=NAVY, bold=True)
bullets(s, 0.7, 5.2, 11.9, 1.7, [
    "팀원이 크롤링한 올리브영·쿠팡 데이터는 실제 평점 기준 부정 비율이 8~17%로 충분히 높았음",
    "3개 플랫폼을 합치고 키워드 추측 대신 실제 평점(1~2점=부정, 4~5점=긍정)으로 재라벨링",
    "무신사 데이터가 워낙 많아서 그대로 합치면 부정 비율이 다시 낮아짐 → 긍정 리뷰 수를 줄여서 부정 25% 비율로 맞춤",
    "최종 학습 데이터: 167,996건 (부정 41,999건 / 긍정 125,997건)",
], size=14.5, color=TEXT, space_after=7)

# ════════════════════════════════════════════════════════════════
# 슬라이드 4 — AI 모델 구현
# ════════════════════════════════════════════════════════════════
s = add_slide(WHITE)
textbox(s, 0.7, 0.5, 11.9, 0.7, "03. AI 모델 구현", size=30, color=NAVY, bold=True)

# 좌측: 전처리 단계 (세로 타임라인)
textbox(s, 0.7, 1.4, 5.6, 0.4, "전처리 과정", size=16, color=NAVY, bold=True)
steps = [
    "결측치 제거",
    "정제 (한글·공백만 남기기)",
    "중복 리뷰 제거",
    "형태소 분석 (Okt 토큰화)",
]
step_y = 1.95
step_h = 0.78
for i, step in enumerate(steps):
    y = step_y + i * step_h
    oval(s, 0.7, y, 0.45, NAVY)
    textbox(s, 0.7, y, 0.45, 0.45, str(i + 1), size=15, color=WHITE, bold=True,
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    rounded_rect(s, 1.35, y, 4.95, 0.55, LIGHT, radius=0.15)
    textbox(s, 1.55, y, 4.6, 0.55, step, size=13.5, color=TEXT, anchor=MSO_ANCHOR.MIDDLE)
textbox(s, 0.7, 5.25, 5.6, 0.5, "※ 16만여 건 형태소 분석에 약 5분 소요",
        size=12, color=GRAY, italic=True)

# 우측: 신경망 구조
textbox(s, 6.9, 1.4, 5.7, 0.4, "신경망 구조 (LSTM, 기준 모델)", size=16, color=NAVY, bold=True)
layers = ["Embedding\n(단어 임베딩)", "LSTM (64)", "Dense (16)\ntanh", "Dense (2)\nsoftmax"]
layer_colors = [NAVY, CORAL, GOLD, GREEN]
lw, lh, lgap = 1.25, 1.6, 0.28
total_lw = lw * 4 + lgap * 3
lx = 6.9 + (5.7 - total_lw) / 2
ly = 2.2
for i, (layer, c) in enumerate(zip(layers, layer_colors)):
    x = lx + i * (lw + lgap)
    rounded_rect(s, x, ly, lw, lh, c, radius=0.1)
    textbox(s, x + 0.08, ly, lw - 0.16, lh, layer, size=12, color=WHITE, bold=True,
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.1)
    if i < 3:
        textbox(s, x + lw, ly + lh / 2 - 0.25, lgap, 0.5, "→", size=20, color=GRAY,
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

textbox(s, 6.9, 4.1, 5.7, 0.8,
        "리뷰는 문장이라 단어의 순서·맥락이 중요 → 순환신경망(LSTM) 사용\n"
        "입력: 정수 인코딩 후 100단어 길이로 패딩",
        size=13, color=TEXT, line_spacing=1.3)

rounded_rect(s, 6.9, 5.1, 5.7, 1.65, NAVY, radius=0.06)
textbox(s, 6.9, 5.22, 5.7, 0.35, "테스트 정확도 (LSTM)", size=13, color=GOLD, bold=True, align=PP_ALIGN.CENTER)
textbox(s, 6.9, 5.58, 5.7, 0.7, "94.2%", size=36, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
textbox(s, 6.9, 6.32, 5.7, 0.35, "부정 정밀도 89.7% · 재현율 86.6%",
        size=11, color=RGBColor(0xC8, 0xCF, 0xE8), align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════════
# 슬라이드 4.5 — 모델 비교 (LSTM vs GRU vs Transformer)
# ════════════════════════════════════════════════════════════════
s = add_slide(WHITE)
textbox(s, 0.7, 0.5, 11.9, 0.7, "모델 비교 — LSTM · GRU · Transformer", size=28, color=NAVY, bold=True)
textbox(s, 0.7, 1.15, 11.9, 0.4,
        "동일한 데이터·전처리로 3가지 신경망 구조를 학습해 비교",
        size=14, color=GRAY, italic=True)

cmp_rows = [
    ("모델", "정확도", "부정 F1", "파라미터", "학습시간", "에포크"),
    ("LSTM", "94.2%", "88.2%", "130만개", "269초", "9"),
    ("GRU", "94.3%", "88.5%", "130만개", "207초", "7"),
    ("Transformer", "93.8%", "87.3%", "129만개", "118초", "7"),
]
ct_x, ct_y, ct_w = 0.9, 1.9, 11.5
col_ws = [2.7, 1.76, 1.76, 1.76, 1.76, 1.76]
row_h2 = 0.62
best_row = 2  # GRU
for ri, row in enumerate(cmp_rows):
    is_header = ri == 0
    is_best = ri == best_row
    bg = NAVY if is_header else (RGBColor(0xE8, 0xF6, 0xF0) if is_best else WHITE)
    rect(s, ct_x, ct_y + ri * row_h2, ct_w, row_h2, bg,
         line=True, line_color=RGBColor(0xE2, 0xE5, 0xEE), line_w=0.75)
    cx = ct_x
    for ci, cell in enumerate(row):
        color = WHITE if is_header else (GREEN if is_best else TEXT)
        textbox(s, cx + 0.2, ct_y + ri * row_h2, col_ws[ci] - 0.3, row_h2, cell,
                size=14, color=color, bold=(is_header or is_best),
                anchor=MSO_ANCHOR.MIDDLE)
        cx += col_ws[ci]

textbox(s, 0.9, ct_y + len(cmp_rows) * row_h2 + 0.25, 11.5, 0.4,
        "GRU가 정확도와 부정 F1 모두 가장 높았고, 학습도 가장 적은 횟수(7회)로 끝났다",
        size=15, color=GREEN, bold=True)

bullets(s, 0.9, ct_y + len(cmp_rows) * row_h2 + 0.75, 11.5, 1.6, [
    "GRU는 구조가 LSTM보다 단순한데도 성능은 비슷하거나 더 좋았음 — 리뷰 문장이 짧아서 충분히 학습 가능",
    "Transformer는 학습 속도가 가장 빨랐지만(118초), 정확도는 셋 중 가장 낮았음",
    "리뷰가 최대 100단어 정도로 짧다 보니, 긴 문맥을 잘 보는 Transformer의 장점이 크게 드러나지 않은 것으로 보임",
    "이번 데이터 규모와 문장 길이에는 GRU가 가장 적합했음",
], size=13.5, color=TEXT, space_after=6)

# ════════════════════════════════════════════════════════════════
# 슬라이드 5 — 배포 및 시연
# ════════════════════════════════════════════════════════════════
s = add_slide(NAVY)
textbox(s, 0.7, 0.5, 11.9, 0.7, "04. 배포 및 시연", size=30, color=WHITE, bold=True)
textbox(s, 0.7, 1.15, 11.9, 0.4, "Streamlit으로 웹 서비스 제작 — 4가지 기능",
        size=14, color=GOLD, italic=True)

features = [
    ("①", "전체 통계", "감성 분포 · 평점 분포\n키워드 Top 20 차트"),
    ("②", "상품별 리뷰 탐색", "상품 선택 → 리뷰 목록\n+ 긍정/부정 비율"),
    ("③", "실시간 감성 예측", "문장을 입력하면\n즉석으로 긍/부정 판단"),
    ("④", "모델 선택", "LSTM·GRU·Transformer\n중 골라서 비교 가능"),
]
fw, fh, fgap = 2.7, 2.5, 0.3
total_fw = fw * 4 + fgap * 3
fx = (W - total_fw) / 2
fy = 1.85
for i, (num, title, desc) in enumerate(features):
    x = fx + i * (fw + fgap)
    rounded_rect(s, x, fy, fw, fh, RGBColor(0x3A, 0x48, 0x8F), radius=0.06)
    oval(s, x + (fw - 0.65) / 2, fy + 0.28, 0.65, CORAL)
    textbox(s, x + (fw - 0.65) / 2, fy + 0.28, 0.65, 0.65, num, size=20, color=WHITE, bold=True,
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    textbox(s, x + 0.15, fy + 1.1, fw - 0.3, 0.42, title, size=14.5, color=WHITE, bold=True,
            align=PP_ALIGN.CENTER)
    textbox(s, x + 0.15, fy + 1.55, fw - 0.3, 0.8, desc, size=11, color=RGBColor(0xCF, 0xD6, 0xEF),
            align=PP_ALIGN.CENTER, line_spacing=1.2)

rounded_rect(s, 1.0, 4.65, 11.33, 2.0, RGBColor(0x24, 0x2F, 0x66), radius=0.06)
textbox(s, 1.3, 4.85, 10.7, 0.4, "시연 영상에서 보여줄 내용", size=14, color=GOLD, bold=True)
textbox(s, 1.3, 5.3, 10.7, 1.2,
        "1. 상품별 리뷰 탐색 페이지에서 리뷰와 감성 비율 확인\n"
        "2. “촉촉하고 발림성도 좋아서 재구매 했어요” 입력 → 긍정\n"
        "3. “피부에 트러블이 생겨서 너무 실망이에요” 입력 → 부정\n"
        "4. 모델을 GRU·Transformer로 바꿔서 같은 문장 다시 예측",
        size=14, color=WHITE, line_spacing=1.35)

prs.save("발표자료.pptx")
print("저장 완료: 발표자료.pptx")
